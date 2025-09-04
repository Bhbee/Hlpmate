import os
import re
import pandas as pd
import docx2txt
import pptx
from typing import List
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from openai import OpenAI
from dotenv import load_dotenv
import streamlit as st

load_dotenv()
client = OpenAI()

class Summarizer:
    @staticmethod
    @st.cache_data(show_spinner=False)
    def extract_text_from_file(file_path: str) -> str:
        ext = file_path.lower().split(".")[-1]
        try:
            if ext == "pdf":
                reader = PdfReader(file_path)
                return "\n".join([page.extract_text() or "" for page in reader.pages])

            elif ext == "docx":
                return docx2txt.process(file_path)

            elif ext == "pptx":
                prs = pptx.Presentation(file_path)
                return "\n".join([
                    shape.text for slide in prs.slides
                    for shape in slide.shapes if hasattr(shape, "text")
                ])

            elif ext == "txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()

            elif ext == "xlsx":
                df = pd.read_excel(file_path)
                return df.astype(str).apply(" ".join, axis=1).str.cat(sep="\n")

            else:
                return ""
        except Exception as e:
            return f"❌ Error reading file: {e}"

    @staticmethod
    def detect_type(text: str) -> str:
        lowered = text.lower()
        if "curriculum vitae" in lowered or "linkedin.com" in lowered:
            return "cv"
        elif "abstract" in lowered and "introduction" in lowered:
            return "academic"
        elif "dear" in lowered and ("invite" in lowered or "welcome" in lowered):
            return "invitation"
        elif "newsletter" in lowered:
            return "newsletter"
        elif "invoice" in lowered or "amount" in lowered:
            return "invoice"
        elif len(lowered.split()) < 200:
            return "short_note"
        else:
            return "generic"

    @staticmethod
    def emphasize_keywords(text: str) -> str:
        keywords = [
            "definition", "meaning", "purpose", "goal", "objective", "summary",
            "key", "important", "exam", "test", "advantages", "disadvantages",
            "merits", "demerits", "characteristics", "features"
        ]
        for word in keywords:
            pattern = r"(?i)\b(" + re.escape(word) + r")\b"
            text = re.sub(pattern, r"**\1**", text)
        return text

    @staticmethod
    def gpt_summarize(prompt: str, max_tokens: int = 300) -> str:
        try:
            response = client.chat.completions.create(
                model="gpt-4-1106-preview",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that summarizes documents clearly and precisely."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.5,
                max_tokens=max_tokens
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            return f"❌ GPT summarization failed: {e}"

    @staticmethod
    @st.cache_data(show_spinner=False)
    def summarize_file(file_path: str) -> str:
        return Summarizer._summarize_file_cached(file_path)

    @staticmethod
    def _summarize_file_cached(file_path: str) -> str:
        full_text = Summarizer.extract_text_from_file(file_path)
        if not full_text or full_text.startswith("❌"):
            return "❌ Could not extract text from the uploaded file."

        doc_type = Summarizer.detect_type(full_text)
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        chunks = splitter.split_text(full_text)

        summaries = []
        for chunk in chunks[:2]:
            prompt = f"Summarize the following {doc_type} document chunk in a clear, useful way for a student:\n\n{chunk}"
            summary = Summarizer.gpt_summarize(prompt)
            summaries.append(summary)

        combined = " ".join(summaries)

        final_prompt = (
            f"Please merge and refine the following summaries from a {doc_type} document "
            f"into a final coherent summary for easy student understanding:\n\n{combined}"
        )
        final_summary = Summarizer.gpt_summarize(final_prompt, max_tokens=600)
        emphasized_summary = Summarizer.emphasize_keywords(final_summary)

        return emphasized_summary
